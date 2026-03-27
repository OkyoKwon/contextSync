"""
ContextSync — PARAIX Hackathon Presentation Generator
Generates a clean, minimal 7-slide PPTX based on 0325PPTX_Plan_contextsync.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ──────────────────────────────────────────────
# Constants
# ──────────────────────────────────────────────

SCREENSHOTS_DIR = os.path.join(os.path.dirname(__file__), "..", "imageAsset", "screenshots")
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "..", "ContextSync_PARAIX_Hackathon.pptx")

# Design tokens
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)
TEXT_SUBTLE = RGBColor(0x64, 0x74, 0x8B)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
DIVIDER_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
OVERLAY_DARK = RGBColor(0x0F, 0x17, 0x2A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_TITLE = "Apple SD Gothic Neo"
FONT_BODY = "Apple SD Gothic Neo"
FONT_MONO = "SF Mono"


def screenshot(name):
    return os.path.join(SCREENSHOTS_DIR, name)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY, line_spacing=1.2):
    """Add a text box and return the text frame for further customization."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(font_size * 0.3)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_paragraph(tf, text, font_size=18, bold=False, color=TEXT_DARK,
                  alignment=PP_ALIGN.LEFT, font_name=FONT_BODY, space_before=0):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    p.space_after = Pt(font_size * 0.3)
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG, line_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_divider(slide, left, top, width, color=DIVIDER_GRAY):
    """Add a thin horizontal line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_badge(slide, left, top, text, bg_color=ACCENT_BLUE, text_color=BG_WHITE, font_size=11):
    """Add a small badge/tag."""
    w = Pt(len(text) * font_size * 0.65 + 24)
    h = Pt(font_size + 14)
    shape = add_rounded_rect(slide, left, top, w, h, fill_color=bg_color)
    shape.adjustments[0] = 0.15
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape


# ──────────────────────────────────────────────
# Slide builders
# ──────────────────────────────────────────────

def build_slide_1_cover(prs):
    """Slide 1: Cover — title + dashboard screenshot"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_WHITE

    # ── Left side: Text content (vertically centered) ──
    left_margin = Inches(0.9)

    # Track badge
    add_badge(slide, left_margin, Inches(2.0), "Track B  |  PARAIX Hackathon",
              bg_color=ACCENT_BLUE, font_size=11)

    # Main title
    add_text_box(slide, left_margin, Inches(2.6), Inches(5.8), Inches(1.0),
                 "ContextSync", font_size=48, bold=True, color=TEXT_DARK)

    # Subtitle
    add_text_box(slide, left_margin, Inches(3.5), Inches(5.8), Inches(0.6),
                 "AI 개발 컨텍스트를 잃지 마세요", font_size=22, color=TEXT_SUBTLE)

    # Tagline
    add_text_box(slide, left_margin, Inches(4.2), Inches(5.8), Inches(0.5),
                 "The Hub for AI Session Context",
                 font_size=15, color=ACCENT_BLUE, font_name=FONT_MONO)

    # Sub-tagline
    add_text_box(slide, left_margin, Inches(4.7), Inches(5.8), Inches(0.5),
                 "세션 관리 · PRD 분석 · 충돌 감지 · AI 평가 — 혼자서도, 팀과 함께도",
                 font_size=13, color=TEXT_SUBTLE)

    # ── Right side: Screenshot ──
    img_path = screenshot("dashboard-full.png")
    if os.path.exists(img_path):
        card = add_rounded_rect(slide, Inches(6.6), Inches(0.7), Inches(6.0), Inches(6.1),
                                fill_color=CARD_BG, line_color=DIVIDER_GRAY)
        slide.shapes.add_picture(img_path, Inches(6.8), Inches(0.9), Inches(5.6))


def build_slide_2_problem(prs):
    """Slide 2: Problem Definition — Before/After comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_WHITE

    # Section label + title
    add_text_box(slide, Inches(0.9), Inches(0.5), Inches(11), Inches(0.4),
                 "문제 정의", font_size=14, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, Inches(0.9), Inches(0.9), Inches(11), Inches(0.7),
                 "Claude Code 팀이 매일 겪는 문제", font_size=28, bold=True, color=TEXT_DARK)

    # Column headers
    before_x = Inches(0.9)
    after_x = Inches(6.7)
    before_w = Inches(5.2)
    after_w = Inches(5.7)
    header_y = Inches(1.8)

    add_text_box(slide, before_x, header_y, before_w, Inches(0.35),
                 "Without ContextSync", font_size=13, bold=True, color=ACCENT_RED)
    add_text_box(slide, after_x, header_y, after_w, Inches(0.35),
                 "With ContextSync", font_size=13, bold=True, color=ACCENT_BLUE)

    # Problem rows
    problems = [
        {
            "before_title": "세션 유실",
            "before_desc": "AI 대화가 로컬 머신에 흩어져 있어\n이전 작업 맥락을 찾을 수 없음",
            "after_title": "중앙 아카이브 + 전문 검색",
            "after_desc": "~/.claude/projects/ 자동 스캔\n웹 대시보드에서 즉시 검색 가능",
            "before_accent": ACCENT_RED,
            "after_accent": ACCENT_BLUE,
            "before_bg": RGBColor(0xFE, 0xF2, 0xF2),
            "after_bg": RGBColor(0xEF, 0xF6, 0xFF),
        },
        {
            "before_title": "충돌 미감지",
            "before_desc": "같은 파일을 동시에 수정해도\n머지할 때까지 아무도 모름",
            "after_title": "실시간 충돌 감지",
            "after_desc": "동시 파일 수정을 자동 탐지\n심각도별 분류 (Info / Warning / Critical)",
            "before_accent": ACCENT_AMBER,
            "after_accent": ACCENT_GREEN,
            "before_bg": RGBColor(0xFF, 0xFB, 0xEB),
            "after_bg": RGBColor(0xEC, 0xFD, 0xF5),
        },
        {
            "before_title": "진척 불투명",
            "before_desc": "팀원이 오늘 뭘 했는지\n파악할 방법이 없음",
            "after_title": "대시보드 + 일간 차트",
            "after_desc": "세션 수, 토큰 사용량, 핫 파일\n7일 활동 통계를 한눈에",
            "before_accent": ACCENT_AMBER,
            "after_accent": ACCENT_BLUE,
            "before_bg": RGBColor(0xFF, 0xFB, 0xEB),
            "after_bg": RGBColor(0xEF, 0xF6, 0xFF),
        },
    ]

    row_h = Inches(1.35)
    row_gap = Inches(0.25)
    start_y = Inches(2.3)
    arrow_w = Inches(0.4)

    for i, prob in enumerate(problems):
        y = start_y + (row_h + row_gap) * i

        # Before card
        card_before = add_rounded_rect(slide, before_x, y, before_w, row_h,
                                       fill_color=prob["before_bg"])
        bar_b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, before_x, y, Pt(5), row_h)
        bar_b.fill.solid()
        bar_b.fill.fore_color.rgb = prob["before_accent"]
        bar_b.line.fill.background()

        add_text_box(slide, before_x + Inches(0.25), y + Inches(0.15),
                     before_w - Inches(0.4), Inches(0.35),
                     prob["before_title"], font_size=16, bold=True,
                     color=prob["before_accent"])
        add_text_box(slide, before_x + Inches(0.25), y + Inches(0.55),
                     before_w - Inches(0.4), Inches(0.65),
                     prob["before_desc"], font_size=13, color=TEXT_SUBTLE, line_spacing=1.4)

        # Arrow
        arrow_x = before_x + before_w + Inches(0.15)
        add_text_box(slide, arrow_x, y + Inches(0.35), arrow_w, Inches(0.5),
                     "→", font_size=24, color=TEXT_SUBTLE, alignment=PP_ALIGN.CENTER)

        # After card
        card_after = add_rounded_rect(slide, after_x, y, after_w, row_h,
                                      fill_color=prob["after_bg"])
        bar_a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, after_x, y, Pt(5), row_h)
        bar_a.fill.solid()
        bar_a.fill.fore_color.rgb = prob["after_accent"]
        bar_a.line.fill.background()

        add_text_box(slide, after_x + Inches(0.25), y + Inches(0.15),
                     after_w - Inches(0.4), Inches(0.35),
                     prob["after_title"], font_size=16, bold=True,
                     color=prob["after_accent"])
        add_text_box(slide, after_x + Inches(0.25), y + Inches(0.55),
                     after_w - Inches(0.4), Inches(0.65),
                     prob["after_desc"], font_size=13, color=TEXT_SUBTLE, line_spacing=1.4)

    # Bottom note
    add_text_box(slide, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
                 "6일간 212커밋 개발 과정에서 직접 겪은 문제입니다",
                 font_size=13, color=TEXT_SUBTLE, alignment=PP_ALIGN.CENTER)


def build_slide_3_solution(prs):
    """Slide 3: Three core features"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_WHITE

    # Section title
    add_text_box(slide, Inches(0.9), Inches(0.6), Inches(11), Inches(0.7),
                 "핵심 기능", font_size=32, bold=True, color=TEXT_DARK)
    add_text_box(slide, Inches(0.9), Inches(1.2), Inches(11), Inches(0.5),
                 "Claude Code 세션을 자동 수집 · 분석 · 공유하는 AI 개발 컨텍스트 허브",
                 font_size=17, color=TEXT_SUBTLE)

    # ── Three feature cards ──
    features = [
        {
            "title": "세션 아카이브 & 검색",
            "desc": "세션을 팀의 지식 자산으로 —\nClaude Code 세션을 자동 수집하고 아카이브합니다.",
            "detail": "~/.claude/projects/ 백그라운드 자동 스캔\nPostgreSQL tsvector 전문 검색\n프로젝트별 그룹핑 · 토큰 비용 분석",
            "color": ACCENT_BLUE,
            "icon": "01",
            "screenshot": "session-conversation.png",
        },
        {
            "title": "실시간 충돌 감지",
            "desc": "머지 전에 충돌을 발견 —\n동시 파일 수정을 실시간으로 탐지합니다.",
            "detail": "심각도 자동 분류 (Critical / Warning / Info)\n파일별 충돌 히트맵 · 리뷰 워크플로우\n팀원 알림 연동",
            "color": ACCENT_AMBER,
            "icon": "02",
            "screenshot": "conflicts-list.png",
        },
        {
            "title": "대시보드 & 분석",
            "desc": "팀 생산성을 한눈에 —\n일간 사용 차트와 7일 활동 통계를 제공합니다.",
            "detail": "세션 수 · 토큰 사용량 · 핫 파일 추적\nPRD 달성률 AI 분석 · AI 활용도 평가\n3가지 배포 모드 (Personal / Team / Member)",
            "color": ACCENT_GREEN,
            "icon": "03",
            "screenshot": "dashboard-stats.png",
        },
    ]

    card_w = Inches(3.7)
    card_h = Inches(4.7)
    start_x = Inches(0.9)
    gap = Inches(0.25)
    card_y = Inches(2.0)

    for i, feat in enumerate(features):
        x = start_x + (card_w + gap) * i

        # Card background
        card = add_rounded_rect(slide, x, card_y, card_w, card_h, fill_color=CARD_BG)

        # Top accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_w, Pt(4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = feat["color"]
        bar.line.fill.background()

        # Number badge
        add_text_box(slide, x + Inches(0.2), card_y + Inches(0.25), Inches(0.5), Inches(0.4),
                     feat["icon"], font_size=14, bold=True, color=feat["color"],
                     font_name=FONT_MONO)

        # Feature title
        add_text_box(slide, x + Inches(0.2), card_y + Inches(0.55), card_w - Inches(0.4), Inches(0.45),
                     feat["title"], font_size=20, bold=True, color=TEXT_DARK)

        # Feature description
        add_text_box(slide, x + Inches(0.2), card_y + Inches(1.0), card_w - Inches(0.4), Inches(0.7),
                     feat["desc"], font_size=13, color=TEXT_SUBTLE, line_spacing=1.4)

        # Detail
        add_text_box(slide, x + Inches(0.2), card_y + Inches(1.7), card_w - Inches(0.4), Inches(0.6),
                     feat["detail"], font_size=11, color=TEXT_SUBTLE, line_spacing=1.4)

        # Screenshot thumbnail
        img_path = screenshot(feat["screenshot"])
        if os.path.exists(img_path):
            img_y = card_y + Inches(2.5)
            slide.shapes.add_picture(img_path, x + Inches(0.2), img_y, card_w - Inches(0.4))

    # Bottom tech stack bar
    add_divider(slide, Inches(0.9), Inches(6.9), Inches(11.5))
    add_text_box(slide, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.4),
                 "React 19  ·  Fastify 5  ·  PostgreSQL 16  ·  Kysely  ·  Claude API  ·  Turborepo",
                 font_size=13, color=TEXT_SUBTLE, alignment=PP_ALIGN.CENTER)


def build_slide_4_ax_4d(prs):
    """Slide 4: AX Know-how — 4D Framework"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_WHITE

    # Title
    add_text_box(slide, Inches(0.9), Inches(0.5), Inches(11), Inches(0.5),
                 "AX 전환 노하우", font_size=14, bold=True, color=ACCENT_BLUE)
    tf = add_text_box(slide, Inches(0.9), Inches(0.9), Inches(11), Inches(0.7),
                      "Claude Code로 Claude Code 관리 도구를 만들다",
                      font_size=28, bold=True, color=TEXT_DARK)

    add_text_box(slide, Inches(0.9), Inches(1.55), Inches(11), Inches(0.4),
                 "개발 전 과정(설계 → 구현 → 테스트 → 문서)을 AI와 협업  |  사용 도구: Claude Code (Opus 4 / Sonnet 4)",
                 font_size=14, color=TEXT_SUBTLE)

    # ── 4D Framework: 4 horizontal cards ──
    phases = [
        {
            "phase": "Discover",
            "title": "탐색",
            "desc": "Claude Code로 기존\n오픈소스 조사 및\n아키텍처 패턴 탐색",
            "color": ACCENT_BLUE,
        },
        {
            "phase": "Define",
            "title": "정의",
            "desc": "PRD · 설계문서를\nClaude가 초안 생성,\n사람이 검토 · 수정",
            "color": RGBColor(0x7C, 0x3A, 0xED),  # violet
        },
        {
            "phase": "Develop",
            "title": "개발",
            "desc": "14개 API 모듈 구현\n27개 DB 마이그레이션\nTDD 사이클",
            "color": ACCENT_GREEN,
        },
        {
            "phase": "Deliver",
            "title": "배포",
            "desc": "CI/CD 파이프라인\n테스트 커버리지\n문서 자동 동기화",
            "color": ACCENT_AMBER,
        },
    ]

    card_w = Inches(2.75)
    card_h = Inches(2.5)
    start_x = Inches(0.9)
    gap = Inches(0.25)
    card_y = Inches(2.3)

    for i, phase in enumerate(phases):
        x = start_x + (card_w + gap) * i

        card = add_rounded_rect(slide, x, card_y, card_w, card_h, fill_color=CARD_BG)

        # Top accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_w, Pt(4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = phase["color"]
        bar.line.fill.background()

        # Phase name (English)
        add_text_box(slide, x + Inches(0.2), card_y + Inches(0.2), card_w - Inches(0.4), Inches(0.35),
                     phase["phase"], font_size=13, bold=True, color=phase["color"],
                     font_name=FONT_MONO)

        # Phase title (Korean)
        add_text_box(slide, x + Inches(0.2), card_y + Inches(0.55), card_w - Inches(0.4), Inches(0.35),
                     phase["title"], font_size=20, bold=True, color=TEXT_DARK)

        # Description
        add_text_box(slide, x + Inches(0.2), card_y + Inches(1.0), card_w - Inches(0.4), Inches(1.2),
                     phase["desc"], font_size=13, color=TEXT_SUBTLE, line_spacing=1.5)

        # Arrow connector (except last)
        if i < 3:
            arrow_x = x + card_w + Inches(0.02)
            add_text_box(slide, arrow_x, card_y + Inches(1.0), Inches(0.22), Inches(0.4),
                         "→", font_size=20, color=DIVIDER_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Bottom: Agent pipeline ──
    pipeline_y = Inches(5.2)
    add_divider(slide, Inches(0.9), pipeline_y, Inches(11.5))

    add_text_box(slide, Inches(0.9), pipeline_y + Inches(0.2), Inches(11.5), Inches(0.35),
                 "멀티 에이전트 파이프라인", font_size=14, bold=True, color=TEXT_DARK)

    agents = ["planner", "tdd-guide", "code-reviewer", "security-reviewer"]
    agent_w = Inches(2.3)
    agent_h = Inches(0.5)
    agent_gap = Inches(0.55)
    agent_y = pipeline_y + Inches(0.7)

    for i, agent in enumerate(agents):
        x = Inches(0.9) + (agent_w + agent_gap) * i
        shape = add_rounded_rect(slide, x, agent_y, agent_w, agent_h,
                                 fill_color=CARD_BG, line_color=DIVIDER_GRAY)
        shape.adjustments[0] = 0.15
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = agent
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT_BLUE
        p.font.name = FONT_MONO
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if i < 3:
            arrow_x = x + agent_w + Inches(0.1)
            add_text_box(slide, arrow_x, agent_y, Inches(0.35), agent_h,
                         "→", font_size=16, color=TEXT_SUBTLE, alignment=PP_ALIGN.CENTER)

    # Context note
    add_text_box(slide, Inches(0.9), agent_y + Inches(0.65), Inches(11.5), Inches(0.35),
                 "CLAUDE.md + memory 시스템으로 프로젝트 컨텍스트 유지  ·  대화별 자동 세션 아카이빙",
                 font_size=12, color=TEXT_SUBTLE)


def build_slide_5_metrics(prs):
    """Slide 5: Quantitative Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_WHITE

    # Title
    add_text_box(slide, Inches(0.9), Inches(0.5), Inches(11), Inches(0.5),
                 "정량 성과", font_size=14, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, Inches(0.9), Inches(0.9), Inches(11), Inches(0.7),
                 "1인 개발자가 Claude Code와 6일 만에",
                 font_size=28, bold=True, color=TEXT_DARK)
    add_text_box(slide, Inches(0.9), Inches(1.45), Inches(11), Inches(0.7),
                 "풀스택 플랫폼을 완성했습니다",
                 font_size=28, bold=True, color=ACCENT_BLUE)

    # ── Metric grid: 3 rows x 4 cols ──
    metrics = [
        ("6", "일", "개발 기간", ACCENT_BLUE),
        ("212", "회", "총 커밋", ACCENT_GREEN),
        ("242K", "", "Lines of Code", ACCENT_BLUE),
        ("14", "개", "API 모듈", ACCENT_GREEN),
        ("27", "개", "DB 마이그레이션", ACCENT_AMBER),
        ("16", "개", "DB 테이블", ACCENT_BLUE),
        ("24", "개", "UI 컴포넌트", ACCENT_GREEN),
        ("178", "개", "E2E 테스트 케이스", ACCENT_AMBER),
        ("3", "종", "배포 모드", ACCENT_BLUE),
        ("4", "종", "기술 문서", ACCENT_GREEN),
        ("80%+", "", "Branch 커버리지", ACCENT_AMBER),
        ("MIT", "", "오픈소스 라이선스", ACCENT_BLUE),
    ]

    grid_start_x = Inches(0.9)
    grid_start_y = Inches(2.3)
    cell_w = Inches(2.85)
    cell_h = Inches(1.3)
    gap_x = Inches(0.2)
    gap_y = Inches(0.15)

    for i, (num, unit, label, color) in enumerate(metrics):
        row = i // 4
        col = i % 4
        x = grid_start_x + (cell_w + gap_x) * col
        y = grid_start_y + (cell_h + gap_y) * row

        card = add_rounded_rect(slide, x, y, cell_w, cell_h, fill_color=CARD_BG)

        # Big number
        num_text = f"{num}{unit}" if unit else num
        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), cell_w - Inches(0.4), Inches(0.6),
                     num_text, font_size=30, bold=True, color=color, font_name=FONT_MONO)

        # Label
        add_text_box(slide, x + Inches(0.2), y + Inches(0.8), cell_w - Inches(0.4), Inches(0.35),
                     label, font_size=13, color=TEXT_SUBTLE)

    # Bottom note
    add_text_box(slide, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
                 "2026-03-19 ~ 03-25  ·  TypeScript 5.7 strict  ·  pnpm + Turborepo 캐싱  ·  CI/CD GitHub Actions  ·  Zod 런타임 검증",
                 font_size=13, color=TEXT_SUBTLE, alignment=PP_ALIGN.CENTER)


def build_slide_6_more_features(prs):
    """Slide 6: Additional features with screenshots"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_WHITE

    # Title
    add_text_box(slide, Inches(0.9), Inches(0.5), Inches(11), Inches(0.5),
                 "더 많은 기능", font_size=14, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, Inches(0.9), Inches(0.9), Inches(11), Inches(0.7),
                 "데모에서 보여드리지 못한 기능들도 모두 작동합니다",
                 font_size=24, bold=True, color=TEXT_DARK)

    # ── Four screenshot cards (2x2 grid) ──
    cards = [
        {
            "title": "PRD 분석",
            "desc": "PRD 업로드 → Claude가 요구사항 추출\n달성률 자동 계산 · 트렌드 리포트",
            "screenshot": "prd-analysis.png",
            "color": ACCENT_BLUE,
        },
        {
            "title": "AI 활용도 평가",
            "desc": "다차원 AI 활용 점수 · 숙련도 등급\n증거 기반 평가 (메시지 인용)",
            "screenshot": "ai-evaluation.png",
            "color": ACCENT_GREEN,
        },
        {
            "title": "전문 검색",
            "desc": "세션 · 메시지 · 파일 경로 · 태그\nPostgreSQL tsvector 기반 즉시 검색",
            "screenshot": "search-overlay.png",
            "color": ACCENT_AMBER,
        },
        {
            "title": "팀 협업 & 설정",
            "desc": "역할 기반 접근 제어 (Owner / Member)\nDual-DB 라우팅 · Join Code 팀 참여",
            "screenshot": "settings-team.png",
            "color": RGBColor(0x7C, 0x3A, 0xED),
        },
    ]

    card_w = Inches(5.8)
    card_h = Inches(2.45)
    start_x = Inches(0.9)
    gap_x = Inches(0.25)
    gap_y = Inches(0.2)
    start_y = Inches(1.7)

    for i, card_info in enumerate(cards):
        row = i // 2
        col = i % 2
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row

        card = add_rounded_rect(slide, x, y, card_w, card_h, fill_color=CARD_BG)

        # Accent bar (left side)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(5), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = card_info["color"]
        bar.line.fill.background()

        # Title
        add_text_box(slide, x + Inches(0.25), y + Inches(0.15), Inches(2.5), Inches(0.35),
                     card_info["title"], font_size=16, bold=True, color=TEXT_DARK)

        # Description
        add_text_box(slide, x + Inches(0.25), y + Inches(0.5), Inches(2.5), Inches(0.7),
                     card_info["desc"], font_size=11, color=TEXT_SUBTLE, line_spacing=1.4)

        # Screenshot (right side of card)
        img_path = screenshot(card_info["screenshot"])
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, x + Inches(2.9), y + Inches(0.15),
                                     Inches(2.7))

    # Bottom note
    add_text_box(slide, Inches(0.9), Inches(6.9), Inches(11.5), Inches(0.4),
                 "이름만 입력하면 바로 사용  ·  Personal / Team Host / Team Member 3가지 배포 모드  ·  Zero-Docker 팀 합류",
                 font_size=13, color=TEXT_SUBTLE, alignment=PP_ALIGN.CENTER)


def build_slide_7_closing(prs):
    """Slide 7: Closing — one-liner + immediate applicability"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = OVERLAY_DARK

    # Center content
    center_y = Inches(1.8)

    # Main message
    add_text_box(slide, Inches(1.5), center_y, Inches(10.3), Inches(1.0),
                 "AI로 개발하는 시대,", font_size=36, bold=True,
                 color=RGBColor(0xF8, 0xFA, 0xFC), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), center_y + Inches(0.8), Inches(10.3), Inches(1.0),
                 "개발 컨텍스트도 AI로 관리하자", font_size=36, bold=True,
                 color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # Divider
    div_y = center_y + Inches(2.0)
    add_divider(slide, Inches(4.5), div_y, Inches(4.3), color=RGBColor(0x33, 0x44, 0x55))

    # Setup command
    cmd_y = div_y + Inches(0.5)
    cmd_bg = add_rounded_rect(slide, Inches(3.2), cmd_y, Inches(6.9), Inches(0.7),
                               fill_color=RGBColor(0x1E, 0x29, 0x3B))
    cmd_bg.adjustments[0] = 0.1
    add_text_box(slide, Inches(3.4), cmd_y + Inches(0.1), Inches(6.5), Inches(0.5),
                 "bash scripts/setup.sh && pnpm dev", font_size=18,
                 color=ACCENT_GREEN, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)

    # Bottom info
    info_y = cmd_y + Inches(1.0)
    features_text = "즉시 적용 가능  ·  오픈소스 (MIT)  ·  파라메타 개발팀 전원 사용 가능"
    add_text_box(slide, Inches(1.5), info_y, Inches(10.3), Inches(0.5),
                 features_text, font_size=15,
                 color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)

    # Project name at bottom
    add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.4),
                 "ContextSync", font_size=16, bold=True,
                 color=RGBColor(0x64, 0x74, 0x8B), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(6.55), Inches(10.3), Inches(0.35),
                 "Claude Code 팀을 위한 세션 아카이브 & 충돌 감지 허브",
                 font_size=12, color=RGBColor(0x47, 0x55, 0x69), alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1_cover(prs)
    build_slide_2_problem(prs)
    build_slide_3_solution(prs)
    build_slide_4_ax_4d(prs)
    build_slide_5_metrics(prs)
    build_slide_6_more_features(prs)
    build_slide_7_closing(prs)

    prs.save(OUTPUT_PATH)
    print(f"Presentation saved: {OUTPUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
